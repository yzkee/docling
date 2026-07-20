from __future__ import annotations

import logging
import shutil
import warnings
from io import BytesIO
from pathlib import Path
from tempfile import mkdtemp
from typing import Any, Callable, Final, Optional, Union

from docling_core.types.doc import (
    BoundingBox,
    CoordOrigin,
    DocItemLabel,
    DoclingDocument,
    DocumentOrigin,
    GroupLabel,
    ImageRef,
    PictureClassificationLabel,
    PictureClassificationMetaField,
    PictureClassificationPrediction,
    PictureMeta,
    ProvenanceItem,
    Size,
    TableCell,
    TableData,
    TabularChartMetaField,
)
from docling_core.types.doc.document import ContentLayer
from lxml import etree
from PIL import Image, UnidentifiedImageError
from typing_extensions import override

from docling.backend.abstract_backend import (
    DeclarativeDocumentBackend,
    PaginatedDocumentBackend,
)
from docling.backend.docx.drawingml.utils import convert_to_modern_format
from docling.datamodel.backend_options import MsPowerpointBackendOptions
from docling.datamodel.base_models import FormatToMimeType, InputFormat
from docling.datamodel.document import InputDocument
from docling.exceptions import DocumentLoadError

_log = logging.getLogger(__name__)

_PPTX_AVAILABLE: bool = False
_PPTX_IMPORT_ERROR: ImportError | None = None
try:  # pragma: no cover - import-time guard
    from pptx import Presentation, presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    from pptx.exc import InvalidXmlError
    from pptx.oxml.text import CT_TextLineBreak

    _PPTX_AVAILABLE = True
except ImportError as e:  # pragma: no cover - import-time guard
    _PPTX_IMPORT_ERROR = e

# Chart image rendering is opt-in and relies on pypdfium2 plus the shared
# DrawingML/LibreOffice helpers, which live behind the PDF extra rather than
# format-pptx. Guard them separately so a slim PPTX install still parses text,
# tables, and chart data; only render_chart_images needs these.
_CHART_RENDER_AVAILABLE: bool = False
try:  # pragma: no cover - import-time guard
    import pypdfium2

    from docling.backend.docx.drawingml.utils import (
        crop_whitespace,
        get_docx_to_pdf_converter,
    )

    _CHART_RENDER_AVAILABLE = True
except ImportError:  # pragma: no cover - import-time guard
    pass

_INSTALL_HINT = (
    "The 'python-pptx' package is required to process PowerPoint files. "
    "Install it with `pip install 'docling-slim[format-pptx]'`."
)

_CHART_RENDER_HINT = (
    "LibreOffice is required to render PowerPoint charts as images "
    "(render_chart_images=True). Install LibreOffice and make sure `soffice` is "
    "on PATH. Charts still keep their classification and reconstructed tabular "
    "data without it."
)

_SAFE_XML_PARSER: Final = etree.XMLParser(
    resolve_entities=False,
    load_dtd=False,
    no_network=True,
    dtd_validation=False,
)
"""Safe XML parser to prevent XXE, DTD-over-network and entity-expansion attacks."""


class MsPowerpointDocumentBackend(DeclarativeDocumentBackend, PaginatedDocumentBackend):
    """Backend for parsing PowerPoint presentations (PPTX and PPT files).

    Converts presentations into structured DoclingDocument format, extracting
    text, tables, images, lists, and comments from slides.

    Legacy ``.ppt`` files (binary PowerPoint 97-2003 format) are first converted
    to ``.pptx`` via LibreOffice before parsing.

    Note:
        Comments are extracted and added to the NOTES content layer. Unlike Word documents,
            PPTX comments only contain position coordinates (x, y) and do not directly reference
            specific shapes or text ranges. Therefore, comments are associated with their parent
            slide group rather than specific document elements.
    """

    # XML namespaces for element lookup
    NAMESPACES: Final[dict[str, str]] = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    # XML relationship types
    COMMENT_REL = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
    )
    COMMENT_AUTHORS_REL = (
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
        "commentAuthors"
    )

    def __init__(
        self,
        in_doc: InputDocument,
        path_or_stream: Union[BytesIO, Path],
        options: Optional[MsPowerpointBackendOptions] = None,
    ) -> None:
        if not _PPTX_AVAILABLE:
            raise ImportError(_INSTALL_HINT) from _PPTX_IMPORT_ERROR
        if in_doc.format == InputFormat.PPT:
            path_or_stream = convert_to_modern_format(path_or_stream, "ppt", "pptx")
        if options is None:
            options = MsPowerpointBackendOptions()
        super().__init__(in_doc, path_or_stream, options)
        self.path_or_stream: Union[BytesIO, Path] = path_or_stream
        self.page_range = in_doc.limits.page_range

        self.pptx_to_pdf_converter: Optional[Callable] = None
        self.pptx_to_pdf_converter_init: bool = False
        self._render_charts: bool = False

        self.pptx_obj: Optional[presentation.Presentation] = None
        self.valid: bool = False
        try:
            if isinstance(self.path_or_stream, BytesIO):
                self.pptx_obj = Presentation(self.path_or_stream)
            elif isinstance(self.path_or_stream, Path):
                self.pptx_obj = Presentation(str(self.path_or_stream))

            self.valid = True
        except Exception as e:
            raise DocumentLoadError(
                f"MsPowerpointDocumentBackend could not load document with hash {self.document_hash}"
            ) from e

        return

    def page_count(self) -> int:
        if self.is_valid():
            assert self.pptx_obj is not None
            return len(self.pptx_obj.slides)
        else:
            return 0

    @override
    def is_valid(self) -> bool:
        return self.valid

    @classmethod
    @override
    def supports_pagination(cls) -> bool:
        return True

    @override
    def unload(self):
        if isinstance(self.path_or_stream, BytesIO):
            self.path_or_stream.close()

        self.path_or_stream = None

    @classmethod
    @override
    def supported_formats(cls) -> set[InputFormat]:
        return {InputFormat.PPTX, InputFormat.PPT}

    @override
    def convert(self) -> DoclingDocument:
        """Parse the PPTX into a structured document model.

        Returns:
            The parsed document.
        """
        origin = DocumentOrigin(
            filename=self.file.name or "file",
            mimetype=FormatToMimeType[self.input_format][0],
            binary_hash=self.document_hash,
        )

        doc = DoclingDocument(name=self.file.stem or "file", origin=origin)
        if self.pptx_obj:
            # Build author map once for all comments
            author_map = self._build_comment_author_map(self.pptx_obj)

            start_page, end_page = self.page_range
            doc = self._walk_linear(
                self.pptx_obj, doc, author_map, start_page=start_page, end_page=end_page
            )

        return doc

    def _generate_prov(
        self, shape, slide_ind, text="", slide_size=Size(width=1, height=1)
    ):
        if shape.left:
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
        else:
            left = 0
            top = 0
            width = slide_size.width
            height = slide_size.height
        shape_bbox = [left, top, left + width, top + height]
        shape_bbox = BoundingBox.from_tuple(shape_bbox, origin=CoordOrigin.BOTTOMLEFT)
        prov = ProvenanceItem(
            page_no=slide_ind + 1, charspan=[0, len(text)], bbox=shape_bbox
        )

        return prov

    def _get_paragraph_level(self, paragraph) -> int:
        """Return the indentation level of a paragraph XML element.

        Paragraphs can have different indentation levels (0-8). The level is
        stored in the `lvl` attribute of the `a:pPr` element (paragraph properties).

        Args:
            paragraph: Paragraph XML element whose level should be extracted.

        Returns:
            Paragraph level in the range (0, 8). Returns 0 when no `a:pPr` element is
                found, no `lvl` attribute exists, or the `lvl` attribute value is
                invalid.
        """
        pPr = paragraph.find("a:pPr", namespaces=self.NAMESPACES)
        if pPr is not None and "lvl" in pPr.attrib:
            try:
                return int(pPr.get("lvl"))
            except ValueError:
                pass
        return 0

    def _parse_bullet_from_paragraph_properties(
        self, pPr
    ) -> tuple[Optional[bool], Optional[str], Optional[str]]:
        """Parse bullet or numbering information from a paragraph properties node.

        This inspects the `a:pPr` or `a:lvlXpPr` element and extracts
        information about the bullet character, automatic numbering, picture
        bullets, or explicit `buNone` markers.

        Args:
            pPr: Paragraph properties XML element (`a:pPr` or `a:lvlXpPr`).

        Returns:
            A 3-tuple (`is_list`, `kind`, `detail`) where: `is_list` is True/False/None
                indicating whether this is a list item; `kind` is one of `buChar`,
                `buAutoNum`, `buBlip`, `buNone` or None, describing the marker type;
                `detail` is the bullet character, numbering type string, or None if not
                applicable.
        """
        if pPr is None:
            return (None, None, None)

        # Explicitly no bullet
        if pPr.find("a:buNone", namespaces=self.NAMESPACES) is not None:
            return (False, "buNone", None)

        # Bullet character
        buChar = pPr.find("a:buChar", namespaces=self.NAMESPACES)
        if buChar is not None:
            return (True, "buChar", buChar.get("char"))

        # Auto numbering
        buAuto = pPr.find("a:buAutoNum", namespaces=self.NAMESPACES)
        if buAuto is not None:
            return (True, "buAutoNum", buAuto.get("type"))

        # Picture bullet
        buBlip = pPr.find("a:buBlip", namespaces=self.NAMESPACES)
        if buBlip is not None:
            return (True, "buBlip", "image")

        return (None, None, None)

    def _find_level_properties_in_list_style(self, lstStyle, lvl: int):
        """Find the level-specific paragraph properties node from a list style.

        This looks for an `a:lvl{lvl+1}pPr` node inside an `a:lstStyle` element, where
        `a:lvl1pPr` corresponds to level 0, `a:lvl2pPr` to level 1, and so on.

        Args:
            lstStyle: List style XML element `a:lstStyle`.
            lvl: Paragraph level in the range (0, 8).

        Returns:
            Matching `a:lvl{lvl+1}pPr` XML element, or None if no matching element is
                found.
        """
        if lstStyle is None:
            return None
        tag = f"a:lvl{lvl + 1}pPr"
        return lstStyle.find(tag, namespaces=self.NAMESPACES)

    def _parse_bullet_from_text_body_list_style(
        self, txBody, lvl: int
    ) -> tuple[Optional[bool], Optional[str], Optional[str]]:
        """Parse bullet or numbering information from a text body's list style.

        This searches for `a:lstStyle/a:lvl{lvl+1}pPr` under a `txBody` and uses the
        level-specific paragraph properties to deduce bullet or numbering information.

        Args:
            txBody: Text body XML element `p:txBody`.
            lvl: Paragraph level in the range (0, 8).

        Returns:
            A 3-tuple (`is_list`, `kind`, `detail`) where: `is_list` is True/False/None
                indicating whether this is a list item; `kind` is one of `buChar`,
                `buAutoNum`, `buBlip`, `buNone` or None, describing the marker type;
                `detail` is the bullet character, numbering type string, or None if not
                applicable.
        """
        if txBody is None:
            return (None, None, None)
        lstStyle = txBody.find("a:lstStyle", namespaces=self.NAMESPACES)
        lvl_pPr = self._find_level_properties_in_list_style(lstStyle, lvl)
        is_list, kind, detail = self._parse_bullet_from_paragraph_properties(lvl_pPr)
        return (is_list, kind, detail)

    def _get_master_text_style_node(
        self, slide_master, placeholder_type
    ) -> Optional[etree._Element]:
        """Get the appropriate master text style node for a placeholder.

        Most content placeholders (BODY/OBJECT) use `p:bodyStyle`, while titles use
        `p:titleStyle`. All other placeholders default to `p:otherStyle`.

        Args:
            slide_master: Slide master object associated with the current slide.
            placeholder_type: Placeholder type enum from `PP_PLACEHOLDER`.

        Returns:
            Matching style node from master `p:txStyles` (`p:bodyStyle`, `p:titleStyle`
                or `p:otherStyle`) or None when no styles are defined.
        """
        txStyles = slide_master._element.find(
            ".//p:txStyles", namespaces=self.NAMESPACES
        )
        if txStyles is None:
            return None

        if placeholder_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            return txStyles.find("p:bodyStyle", namespaces=self.NAMESPACES)

        if placeholder_type == PP_PLACEHOLDER.TITLE:
            return txStyles.find("p:titleStyle", namespaces=self.NAMESPACES)

        return txStyles.find("p:otherStyle", namespaces=self.NAMESPACES)

    def _parse_bullet_from_master_text_styles(
        self, slide_master, placeholder_type, lvl: int
    ) -> tuple[Optional[bool], Optional[str], Optional[str]]:
        """Parse bullet or numbering information from the slide master text styles.

        This looks up the appropriate style bucket in the slide master's `p:txStyles`
        (`titleStyle`, `bodyStyle` or `otherStyle`) and extracts bullet or numbering
        information for the given level.

        Args:
            slide_master: Slide master object associated with the current slide.
            placeholder_type: Placeholder type enum from `PP_PLACEHOLDER`.
            lvl: Paragraph level in the range (0, 8).

        Returns:
            A 3-tuple (`is_list`, `kind`, `detail`) where: `is_list` is True/False/None
                indicating whether this is a list item; `kind` is one of `buChar`,
                `buAutoNum`, `buBlip`, `buNone` or None, describing the marker type;
                `detail` is the bullet character, numbering type string, or None if not
                applicable.
        """
        style = self._get_master_text_style_node(slide_master, placeholder_type)
        if style is None:
            return (None, None, None)

        lvl_pPr = style.find(f".//a:lvl{lvl + 1}pPr", namespaces=self.NAMESPACES)
        is_list, kind, detail = self._parse_bullet_from_paragraph_properties(lvl_pPr)
        return (is_list, kind, detail)

    def _is_list_item(self, paragraph) -> tuple[bool, str]:
        """Determine whether a paragraph should be treated as a list item.

        The method first tries to resolve list style information via the shape that
        owns the paragraph. If that is not possible, it falls back to simpler checks
        based on paragraph properties and level.

        Args:
            paragraph: `python-pptx` paragraph object to inspect.

        Returns:
            A 2-tuple (`is_list`, `bullet_type`) where: `is_list` is True if the
                paragraph is considered a list item, otherwise False; `bullet_type` is
                one of `Bullet`, `Numbered` or `None`, describing the list marker type.
        """
        p = paragraph._element

        # Try to get shape from paragraph if possible
        shape = None
        try:
            # This path works for python-pptx paragraphs
            # First get the text_frame (paragraph's parent)
            text_frame = paragraph._parent
            # Then get the shape (text_frame's parent)
            shape = text_frame._parent
        except AttributeError:
            pass

        if shape is not None:
            marker_info = self._get_effective_list_marker(shape, paragraph)

            # Check if it's definitely a list item
            if marker_info["is_list"] is True or marker_info["kind"] in (
                "buChar",
                "buAutoNum",
                "buBlip",
            ):
                if marker_info["kind"] == "buChar":
                    return (True, "Bullet")
                elif marker_info["kind"] == "buAutoNum":
                    return (True, "Numbered")
                else:
                    return (True, "None")

            # Check if it's definitely not a list item
            if marker_info["is_list"] is False:
                return (False, "None")

            # Fallback to paragraph level check
            if paragraph.level > 0:
                return (True, "None")

            return (False, "None")

        # Fallback to simpler check if shape is not available
        if p.find(".//a:buChar", namespaces={"a": self.NAMESPACES["a"]}) is not None:
            return (True, "Bullet")
        elif (
            p.find(".//a:buAutoNum", namespaces={"a": self.NAMESPACES["a"]}) is not None
        ):
            return (True, "Numbered")
        elif paragraph.level > 0:
            # Most likely a sub-list
            return (True, "None")
        else:
            return (False, "None")

    def _get_effective_list_marker(self, shape, paragraph) -> dict:
        """Return a dictionary describing the effective list marker for a paragraph.

        List marker information can come from several sources: direct paragraph
        properties, shape-level list styles, layout placeholders, or slide master text
        styles. This helper resolves all of these layers and returns a unified view of
        the effective marker.

        Args:
            shape: Shape object that contains the paragraph.
            paragraph: `python-pptx` paragraph object to inspect.

        Returns:
            Information about the list marker in a dictionary, where: `is_list` is
                True/False/None indicating if this is a list item; `kind` is one of
                `buChar`, `buAutoNum`, `buBlip`, `buNone` or None, describing the
                marker type; `detail` is the bullet character or numbering type string,
                or None if not applicable; `level` is the paragraph level in the range
                (0, 8).
        """
        p = paragraph._element
        lvl = self._get_paragraph_level(p)

        # 1) Direct paragraph properties
        pPr = p.find("a:pPr", namespaces=self.NAMESPACES)
        is_list, kind, detail = self._parse_bullet_from_paragraph_properties(pPr)
        if is_list is not None:
            return {
                "is_list": is_list,
                "kind": kind,
                "detail": detail,
                "level": lvl,
            }

        # 2) Shape-level lstStyle (txBody/a:lstStyle)
        txBody = shape._element.find(".//p:txBody", namespaces=self.NAMESPACES)
        is_list, kind, detail = self._parse_bullet_from_text_body_list_style(
            txBody, lvl
        )
        if is_list is not None:
            return {
                "is_list": is_list,
                "kind": kind,
                "detail": detail,
                "level": lvl,
            }

        # 3) Layout placeholder lstStyle (if this is a placeholder)
        layout_result = None
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            layout = shape.part.slide.slide_layout
            layout_ph = None
            try:
                layout_ph = layout.placeholders.get(idx)
            except Exception:
                layout_ph = None

            if layout_ph is not None:
                layout_tx = layout_ph._element.find(
                    ".//p:txBody", namespaces=self.NAMESPACES
                )
                is_list, kind, detail = self._parse_bullet_from_text_body_list_style(
                    layout_tx, lvl
                )

                # Only use layout result if is_list is explicitly True/False
                if is_list is not None:
                    layout_result = {
                        "is_list": is_list,
                        "kind": kind,
                        "detail": detail,
                        "level": lvl,
                    }

                # 4) Parse master txStyles
                ph_type = shape.placeholder_format.type
                master = shape.part.slide.slide_layout.slide_master
                is_list, kind, detail = self._parse_bullet_from_master_text_styles(
                    master, ph_type, lvl
                )

                # Check if master has marker information
                if kind in ("buChar", "buAutoNum", "buBlip"):
                    return {
                        "is_list": True,
                        "kind": kind,
                        "detail": detail,
                        "level": lvl,
                    }
                elif is_list is not None:
                    return {
                        "is_list": is_list,
                        "kind": kind,
                        "detail": detail,
                        "level": lvl,
                    }

            # If layout has explicit is_list value but master didn't override it, use
            # layout
            if layout_result is not None:
                return layout_result

        return {
            "is_list": None,
            "kind": None,
            "detail": None,
            "level": lvl,
        }

    def _handle_text_elements(
        self, shape, parent_slide, slide_ind, doc: DoclingDocument, slide_size
    ):
        is_list_group_created = False
        enum_list_item_value = 0
        new_list = None
        doc_label = DocItemLabel.LIST_ITEM
        prov = self._generate_prov(shape, slide_ind, shape.text.strip(), slide_size)

        # Iterate through paragraphs to build up text
        for paragraph in shape.text_frame.paragraphs:
            is_a_list, bullet_type = self._is_list_item(paragraph)
            p = paragraph._element

            # Convert line breaks to spaces and accumulate text
            p_text = ""
            for e in p.content_children:
                if isinstance(e, CT_TextLineBreak):
                    p_text += " "
                else:
                    p_text += e.text

            if is_a_list:
                enum_marker = ""
                enumerated = bullet_type == "Numbered"

                if not is_list_group_created:
                    new_list = doc.add_list_group(
                        name="list",
                        parent=parent_slide,
                    )
                    is_list_group_created = True
                    enum_list_item_value = 0

                if enumerated:
                    enum_list_item_value += 1
                    enum_marker = str(enum_list_item_value) + "."

                doc.add_list_item(
                    marker=enum_marker,
                    enumerated=enumerated,
                    parent=new_list,
                    text=p_text,
                    prov=prov,
                )
            else:  # is paragraph not a list item
                if is_list_group_created:
                    is_list_group_created = False
                    new_list = None
                    enum_list_item_value = 0
                # Assign proper label to the text, depending if it's a Title or Section Header
                # For other types of text, assign - PARAGRAPH
                doc_label = DocItemLabel.PARAGRAPH
                if shape.is_placeholder:
                    placeholder_type = shape.placeholder_format.type
                    if placeholder_type in [
                        PP_PLACEHOLDER.CENTER_TITLE,
                        PP_PLACEHOLDER.TITLE,
                    ]:
                        # It's a title
                        doc_label = DocItemLabel.TITLE

                # output accumulated inline text:
                doc.add_text(
                    label=doc_label,
                    parent=parent_slide,
                    text=p_text,
                    prov=prov,
                )
        return

    def _handle_title(self, shape, parent_slide, slide_ind, doc):
        placeholder_type = shape.placeholder_format.type
        txt = shape.text.strip()
        prov = self._generate_prov(shape, slide_ind, txt)

        if len(txt.strip()) > 0:
            # title = slide.shapes.title.text if slide.shapes.title else "No title"
            if placeholder_type in [PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.TITLE]:
                _log.info(f"Title found: {shape.text}")
                doc.add_text(
                    label=DocItemLabel.TITLE, parent=parent_slide, text=txt, prov=prov
                )
            elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                _log.info(f"Subtitle found: {shape.text}")
                # Using DocItemLabel.FOOTNOTE, while SUBTITLE label is not avail.
                doc.add_text(
                    label=DocItemLabel.SECTION_HEADER,
                    parent=parent_slide,
                    text=txt,
                    prov=prov,
                )
        return

    def _handle_pictures(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Open it with PIL
        try:
            # Get the image bytes
            image = shape.image
            image_bytes = image.blob
            im_dpi, _ = image.dpi
            pil_image = Image.open(BytesIO(image_bytes))

            # shape has picture
            prov = self._generate_prov(shape, slide_ind, "", slide_size)
            doc.add_picture(
                parent=parent_slide,
                image=ImageRef.from_pil(image=pil_image, dpi=im_dpi),
                caption=None,
                prov=prov,
            )
        except (
            UnidentifiedImageError,
            OSError,
            ValueError,
            InvalidXmlError,
            KeyError,
            AttributeError,
        ) as e:
            warnings.warn(
                f"Skipping malformed picture shape: {e}",
                UserWarning,
                stacklevel=2,
            )
        return

    def _handle_tables(self, shape, parent_slide, slide_ind, doc, slide_size):
        # Handling tables, images, charts
        if shape.has_table:
            table = shape.table
            table_xml = shape._element

            prov = self._generate_prov(shape, slide_ind, "", slide_size)

            num_cols = 0
            num_rows = len(table.rows)
            tcells = []
            # Access the XML element for the shape that contains the table
            table_xml = shape._element

            for row_idx, row in enumerate(table.rows):
                if len(row.cells) > num_cols:
                    num_cols = len(row.cells)
                for col_idx, cell in enumerate(row.cells):
                    # Access the XML of the cell (this is the 'tc' element in table XML)
                    cell_xml = table_xml.xpath(
                        f".//a:tbl/a:tr[{row_idx + 1}]/a:tc[{col_idx + 1}]"
                    )

                    if not cell_xml:
                        continue  # If no cell XML is found, skip

                    cell_xml = cell_xml[0]  # Get the first matching XML node
                    row_span = cell_xml.get("rowSpan")  # Vertical span
                    col_span = cell_xml.get("gridSpan")  # Horizontal span

                    if row_span is None:
                        row_span = 1
                    else:
                        row_span = int(row_span)

                    if col_span is None:
                        col_span = 1
                    else:
                        col_span = int(col_span)

                    icell = TableCell(
                        text=cell.text.strip(),
                        row_span=row_span,
                        col_span=col_span,
                        start_row_offset_idx=row_idx,
                        end_row_offset_idx=row_idx + row_span,
                        start_col_offset_idx=col_idx,
                        end_col_offset_idx=col_idx + col_span,
                        column_header=row_idx == 0,
                        row_header=False,
                    )
                    if len(cell.text.strip()) > 0:
                        tcells.append(icell)
            # Initialize Docling TableData
            data = TableData(num_rows=num_rows, num_cols=num_cols, table_cells=[])
            # Populate
            for tcell in tcells:
                data.table_cells.append(tcell)
            if len(tcells) > 0:
                # If table is not fully empty...
                # Create Docling table
                doc.add_table(parent=parent_slide, data=data, prov=prov)
        return

    @staticmethod
    def _chart_type_to_classification(chart_type: Any) -> PictureClassificationLabel:
        """Map a python-pptx ``XL_CHART_TYPE`` to a docling classification label.

        ``chart.chart_type`` is a granular enum (``COLUMN_CLUSTERED``,
        ``XY_SCATTER_LINES``, ``THREE_D_PIE``, ...). Rather than enumerate every
        member, we match on the enum member's name by family, mirroring the
        Excel backend's tagname mapping: column/bar variants become BAR_CHART,
        line variants LINE_CHART, pie/doughnut PIE_CHART, scatter SCATTER_CHART,
        and everything else (area, radar, stock, surface, ...) OTHER_CHART.
        SCATTER is matched before LINE because ``XY_SCATTER_LINES`` contains both.

        Args:
            chart_type: An ``XL_CHART_TYPE`` member, or None when python-pptx
                cannot determine the type (e.g. combination charts).

        Returns:
            The matching PictureClassificationLabel.
        """
        name = chart_type.name if chart_type is not None else ""
        if "PIE" in name or "DOUGHNUT" in name:
            return PictureClassificationLabel.PIE_CHART
        if "SCATTER" in name:
            return PictureClassificationLabel.SCATTER_CHART
        if "LINE" in name:
            return PictureClassificationLabel.LINE_CHART
        if "BAR" in name or "COL" in name:
            return PictureClassificationLabel.BAR_CHART
        return PictureClassificationLabel.OTHER_CHART

    @staticmethod
    def _chart_title_text(chart: Any) -> Optional[str]:
        """Return the chart's title text, or None when it has no title.

        Args:
            chart: A python-pptx ``Chart`` object.

        Returns:
            The stripped title text, or None.
        """
        if not chart.has_title:
            return None
        text = chart.chart_title.text_frame.text.strip()
        return text or None

    @staticmethod
    def _cell_text(value: Any) -> str:
        """Format a chart category label or data value as cell text.

        python-pptx returns numeric series values as floats and empty points as
        None. Integer-valued floats are rendered without a trailing ``.0`` so the
        reconstructed table reads like the source data (``120`` not ``120.0``).

        Args:
            value: A category label or data point (str, float, int, or None).

        Returns:
            The value as a string ("" for None).
        """
        if value is None:
            return ""
        if isinstance(value, float) and value.is_integer():
            return str(int(value))
        return str(value)

    def _chart_to_table_data(self, chart: Any) -> Optional[TableData]:
        """Reconstruct a chart's underlying data grid as a TableData.

        Layout produced (categories down the first column, one column per series):

            | <blank> | <series 0 name> | <series 1 name> | ...
            | cat_0   | val_0,0         | val_1,0         | ...
            | cat_1   | val_0,1         | val_1,1         | ...

        Unlike the Excel backend, python-pptx exposes the plotted numbers
        directly on each series (``series.values``) and the shared category
        labels on the plot (``plot.categories``) — no workbook reference
        resolution is needed.

        Args:
            chart: A python-pptx ``Chart`` object.

        Returns:
            A TableData, or None if the chart exposes no usable series.
        """
        series_list = list(chart.series)
        if not series_list:
            return None

        plots = list(chart.plots)
        categories: list[str] = (
            [self._cell_text(cat) for cat in plots[0].categories] if plots else []
        )

        columns: list[tuple[str, list[str]]] = []
        for series in series_list:
            name = self._cell_text(series.name) if series.name is not None else ""
            values = [self._cell_text(v) for v in series.values]
            columns.append((name, values))

        num_data_rows = max([len(categories)] + [len(values) for _, values in columns])
        if num_data_rows == 0:
            return None

        num_rows = num_data_rows + 1
        num_cols = 1 + len(columns)
        cells: list[TableCell] = []

        header_labels = [""] + [name for name, _ in columns]
        for col_idx, label in enumerate(header_labels):
            cells.append(
                TableCell(
                    text=label,
                    row_span=1,
                    col_span=1,
                    start_row_offset_idx=0,
                    end_row_offset_idx=1,
                    start_col_offset_idx=col_idx,
                    end_col_offset_idx=col_idx + 1,
                    column_header=True,
                    row_header=False,
                )
            )
        for data_row in range(num_data_rows):
            row_idx = data_row + 1
            category = categories[data_row] if data_row < len(categories) else ""
            row_texts = [category] + [
                (values[data_row] if data_row < len(values) else "")
                for _, values in columns
            ]
            for col_idx, text in enumerate(row_texts):
                cells.append(
                    TableCell(
                        text=text,
                        row_span=1,
                        col_span=1,
                        start_row_offset_idx=row_idx,
                        end_row_offset_idx=row_idx + 1,
                        start_col_offset_idx=col_idx,
                        end_col_offset_idx=col_idx + 1,
                        column_header=False,
                        row_header=(col_idx == 0),
                    )
                )

        return TableData(num_rows=num_rows, num_cols=num_cols, table_cells=cells)

    def _handle_chart(self, shape, parent_slide, slide_ind, doc, slide_size):
        """Add a native chart shape as a classified PictureItem with its data.

        Each chart becomes a PictureItem whose meta carries (a) the chart-type
        classification and (b) the chart's plotted numbers reconstructed as a
        TableData. The chart title, if any, becomes the picture caption. When
        ``render_chart_images`` is enabled and LibreOffice is available, a
        rendered image is attached; on any rendering failure the picture keeps
        its classification and data without an image.

        Args:
            shape: The python-pptx graphic-frame shape holding the chart.
            parent_slide: The parent slide group.
            slide_ind: Zero-based slide index.
            doc: The DoclingDocument to update.
            slide_size: The slide size, used for provenance.
        """
        try:
            chart = shape.chart
        except (ValueError, KeyError, InvalidXmlError) as exc:
            warnings.warn(
                f"Skipping malformed chart shape: {exc}",
                UserWarning,
                stacklevel=2,
            )
            return

        try:
            chart_type = chart.chart_type
        except (ValueError, NotImplementedError):
            chart_type = None
        classification = self._chart_type_to_classification(chart_type)
        caption_text = self._chart_title_text(chart)
        table_data = self._chart_to_table_data(chart)

        prov = self._generate_prov(shape, slide_ind, "", slide_size)

        image_ref = None
        if self._render_charts:
            try:
                chart_image = self._render_chart_image(slide_ind, shape.shape_id)
                if chart_image is not None:
                    image_ref = ImageRef.from_pil(image=chart_image, dpi=72)
            except Exception:
                _log.warning(
                    "could not render a chart image; keeping chart data without image",
                    exc_info=True,
                )

        caption_item = (
            doc.add_text(label=DocItemLabel.CAPTION, text=caption_text)
            if caption_text
            else None
        )

        picture = doc.add_picture(
            parent=parent_slide,
            image=image_ref,
            caption=caption_item,
            prov=prov,
        )
        picture.meta = PictureMeta(
            classification=PictureClassificationMetaField(
                predictions=[PictureClassificationPrediction(class_name=classification)]
            ),
            tabular_chart=(
                TabularChartMetaField(chart_data=table_data)
                if table_data is not None
                else None
            ),
        )
        return

    def _get_libreoffice_converter(self) -> Optional[Callable]:
        """Lazily initialize and return a LibreOffice converter callable.

        The converter accepts ``(input_path, output_path)`` and converts the
        input file to PDF. Returns None when LibreOffice is not available.
        """
        if self.pptx_to_pdf_converter_init:
            return self.pptx_to_pdf_converter

        self.pptx_to_pdf_converter_init = True
        if _CHART_RENDER_AVAILABLE:
            self.pptx_to_pdf_converter = get_docx_to_pdf_converter()
        if self.pptx_to_pdf_converter is None:
            _log.debug("LibreOffice not found — PPTX charts will not be rendered.")
        return self.pptx_to_pdf_converter

    def _isolate_chart_presentation(
        self, slide_ind: int, chart_shape_id: int, out_path: Path
    ) -> bool:
        """Save a copy of the presentation holding only the target chart.

        A fresh copy of the loaded presentation is reopened, every slide except
        the chart's is removed, and on that slide every shape except the chart
        is removed. LibreOffice then renders a single-chart page. When the chart
        is not a top-level shape (e.g. nested in a group) its ``shape_id`` is not
        found among the slide's shapes, so the slide is left intact and the whole
        slide is rendered instead — a best-effort fallback.

        Args:
            slide_ind: Zero-based index of the slide holding the chart.
            chart_shape_id: The ``shape_id`` of the chart's graphic frame.
            out_path: Destination path for the trimmed ``.pptx``.

        Returns:
            True when the trimmed presentation was written.
        """
        if self.pptx_obj is None:
            return False

        buf = BytesIO()
        self.pptx_obj.save(buf)
        buf.seek(0)
        prs = Presentation(buf)

        slide_id_list = prs.slides._sldIdLst
        slide_ids = list(slide_id_list)
        if slide_ind >= len(slide_ids):
            return False
        target_slide = prs.slides[slide_ind]

        for idx, slide_id in enumerate(slide_ids):
            if idx != slide_ind:
                slide_id_list.remove(slide_id)

        for shp in list(target_slide.shapes):
            if shp.shape_id != chart_shape_id:
                shp._element.getparent().remove(shp._element)

        prs.save(str(out_path))
        return True

    def _render_chart_image(
        self, slide_ind: int, chart_shape_id: int
    ) -> Optional[Image.Image]:
        """Render a native chart to an image via LibreOffice.

        PPTX stores charts as vector definitions with no embedded raster. To
        obtain a picture we isolate the chart onto a throwaway single-slide
        presentation, convert that to PDF with LibreOffice — the same external
        tool already used for EMF/WMF images — and rasterize the first page with
        pypdfium2, trimming the surrounding whitespace.

        Args:
            slide_ind: Zero-based index of the slide holding the chart.
            chart_shape_id: The ``shape_id`` of the chart's graphic frame.

        Returns:
            A PIL Image, or None when LibreOffice is unavailable or the
            conversion fails.
        """
        converter = self._get_libreoffice_converter()
        if converter is None:
            return None

        temp_dir = Path(mkdtemp())
        try:
            input_path = temp_dir / "chart.pptx"
            output_path = temp_dir / "chart.pdf"
            if not self._isolate_chart_presentation(
                slide_ind, chart_shape_id, input_path
            ):
                return None
            converter(input_path, output_path)
            if not output_path.exists():
                _log.debug("LibreOffice produced no PDF output for a chart")
                return None
            pdf = pypdfium2.PdfDocument(str(output_path))
            page = pdf[0]
            pil_image = crop_whitespace(page.render(scale=2).to_pil())
            page.close()
            pdf.close()
            return pil_image
        except Exception as exc:
            _log.debug("Chart rendering via LibreOffice failed: %s", exc)
            return None
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def _walk_linear(
        self,
        pptx_obj: presentation.Presentation,
        doc: DoclingDocument,
        author_map: dict[str, tuple[str, str]],
        start_page: int = 1,
        end_page: Optional[int] = None,
    ) -> DoclingDocument:
        # Units of size in PPTX by default are EMU units (English Metric Units)
        slide_width = pptx_obj.slide_width
        slide_height = pptx_obj.slide_height

        self._render_charts = (
            isinstance(self.options, MsPowerpointBackendOptions)
            and self.options.render_chart_images
        )
        if self._render_charts and self._get_libreoffice_converter() is None:
            _log.warning(_CHART_RENDER_HINT)
            self._render_charts = False

        max_levels = 10
        parents = {}  # type: ignore
        for i in range(max_levels):
            parents[i] = None

        # Loop through each slide
        selected_slides = list(enumerate(pptx_obj.slides))[start_page - 1 : end_page]
        for slide_ind, slide in selected_slides:
            parent_slide = doc.add_group(
                name=f"slide-{slide_ind}", label=GroupLabel.CHAPTER, parent=parents[0]
            )

            slide_size = Size(width=slide_width, height=slide_height)
            doc.add_page(page_no=slide_ind + 1, size=slide_size)

            def _safe_shape_type(shape):
                """Return shape.shape_type, or None if unrecognized.

                python-pptx raises NotImplementedError for <p:sp> elements
                that don't match any known shape category (placeholder,
                freeform, autoshape, textbox).
                """
                try:
                    return shape.shape_type
                except NotImplementedError:
                    _log.debug("Skipping shape with unrecognized type: %s", shape.name)
                    return None

            def handle_shapes(shape, parent_slide, slide_ind, doc, slide_size):
                handle_groups(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.has_table:
                    # Handle Tables
                    self._handle_tables(shape, parent_slide, slide_ind, doc, slide_size)
                if shape.has_chart:
                    self._handle_chart(shape, parent_slide, slide_ind, doc, slide_size)
                if _safe_shape_type(shape) == MSO_SHAPE_TYPE.PICTURE:
                    # Handle Pictures
                    self._handle_pictures(
                        shape, parent_slide, slide_ind, doc, slide_size
                    )
                # If shape doesn't have any text, move on to the next shape
                if not hasattr(shape, "text"):
                    return
                if shape.text is None:
                    return
                if len(shape.text.strip()) == 0:
                    return
                if not shape.has_text_frame:
                    _log.warning("Warning: shape has text but not text_frame")
                    return
                # Handle other text elements, including lists (bullet lists, numbered
                # lists)
                self._handle_text_elements(
                    shape, parent_slide, slide_ind, doc, slide_size
                )
                return

            def handle_groups(shape, parent_slide, slide_ind, doc, slide_size):
                if _safe_shape_type(shape) == MSO_SHAPE_TYPE.GROUP:
                    for groupedshape in shape.shapes:
                        handle_shapes(
                            groupedshape, parent_slide, slide_ind, doc, slide_size
                        )

            # Loop through each shape in the slide
            for shape in slide.shapes:
                handle_shapes(shape, parent_slide, slide_ind, doc, slide_size)

            # Handle notes slide
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame is not None:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        bbox = BoundingBox(l=0, t=0, r=0, b=0)
                        prov = ProvenanceItem(
                            page_no=slide_ind + 1,
                            charspan=[0, len(notes_text)],
                            bbox=bbox,
                        )
                        doc.add_text(
                            label=DocItemLabel.TEXT,
                            parent=parent_slide,
                            text=notes_text,
                            prov=prov,
                            content_layer=ContentLayer.NOTES,
                        )

            # Extract comments for this slide
            self._extract_slide_comments(
                slide, slide_ind, author_map, doc, parent_slide
            )

        return doc

    def _build_comment_author_map(
        self, pptx_obj: presentation.Presentation
    ) -> dict[str, tuple[str, str]]:
        """Build a map of comment author IDs to (name, initials).

        Args:
            pptx_obj: The PowerPoint presentation object.

        Returns:
            Dictionary mapping author ID to (name, initials) tuple.
        """
        author_map: dict[str, tuple[str, str]] = {}
        try:
            for rel in pptx_obj.part.rels.values():
                if rel.reltype == self.COMMENT_AUTHORS_REL:
                    root = etree.fromstring(
                        rel.target_part.blob, parser=_SAFE_XML_PARSER
                    )
                    author_map = {
                        author_el.get("id", ""): (
                            author_el.get("name", ""),
                            author_el.get("initials", ""),
                        )
                        for author_el in root.findall(
                            "p:cmAuthor", namespaces=self.NAMESPACES
                        )
                    }
        except Exception as e:
            _log.debug(f"Could not parse PPTX comment authors: {e}")
        return author_map

    def _extract_slide_comments(
        self,
        slide,
        slide_idx: int,
        author_map: dict[str, tuple[str, str]],
        doc: DoclingDocument,
        parent_slide,
    ) -> None:
        """Extract and add comments for a specific slide.

        Args:
            slide: The slide object to extract comments from.
            slide_idx: Zero-based slide index.
            author_map: Dictionary mapping author ID to (name, initials).
            doc: The DoclingDocument to add comments to.
            parent_slide: The parent slide group to link comments to.
        """
        try:
            slide_part = slide.part
        except Exception as e:
            _log.debug(
                f"Could not access slide part for slide {slide_idx + 1}, "
                f"skipping comment extraction: {e}"
            )
            return

        for rel in slide_part.rels.values():
            if rel.reltype != self.COMMENT_REL:
                continue
            try:
                root = etree.fromstring(rel.target_part.blob, parser=_SAFE_XML_PARSER)
                for cm in root.findall("p:cm", namespaces=self.NAMESPACES):
                    author_id = cm.get("authorId", "")
                    dt = cm.get("dt", "")
                    text_el = cm.find("p:text", namespaces=self.NAMESPACES)
                    raw_text = (
                        (text_el.text or "").strip() if text_el is not None else ""
                    )
                    if not raw_text:
                        continue

                    name, initials = author_map.get(author_id, ("", ""))
                    metadata_parts = []
                    if name:
                        author_str = f"author: {name}"
                        if initials:
                            author_str += f" ({initials})"
                        metadata_parts.append(author_str)
                    if dt:
                        metadata_parts.append(f"time: {dt}")
                    prefix = ", ".join(metadata_parts)
                    full_text = f"[{prefix}]: {raw_text}" if prefix else raw_text

                    comment_idx = cm.get("idx", str(slide_idx))
                    comment_group = doc.add_group(
                        label=GroupLabel.COMMENT_SECTION,
                        name=f"comment-slide{slide_idx + 1}-{comment_idx}",
                        content_layer=ContentLayer.NOTES,
                    )
                    doc.add_comment(
                        text=full_text,
                        targets=None,
                        parent=comment_group,
                    )
                    _log.debug(
                        f"Added PPTX comment slide {slide_idx + 1} idx={comment_idx}"
                    )
            except Exception as e:
                _log.debug(f"Could not parse comments for slide {slide_idx}: {e}")
